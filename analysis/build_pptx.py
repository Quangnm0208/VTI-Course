"""
build_pptx.py — Sinh bài thuyết trình 10 slide (.pptx) từ bộ số liệu đã chốt.

Chạy:
    python analysis/build_pptx.py

Đầu ra:
    docs/presentation.pptx  (10 slide, biểu đồ native có thể chỉnh trong PowerPoint,
                             speaker notes đầy đủ)
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

sys.path.insert(0, str(Path(__file__).resolve().parent))
import market_insight_data as D  # noqa: E402

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "docs" / "presentation.pptx"

# Palette
NAVY = RGBColor(0x1E, 0x3A, 0x8A)
BLUE = RGBColor(0x25, 0x63, 0xEB)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
RED = RGBColor(0xDC, 0x26, 0x26)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
INK = RGBColor(0x0F, 0x17, 0x2A)
GREY = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)

W, H = Inches(13.333), Inches(7.5)


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, x, y, w, h, fill=None, line=None):
    sh = slide.shapes.add_shape(1, x, y, w, h)  # rectangle
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill if fill else WHITE
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """runs: list of (text, size, bold, color) hoặc list các paragraph (list runs)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if runs and isinstance(runs[0], tuple):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for (txt, size, bold, color) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Segoe UI"
    return tb


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _header(slide, idx, title, pct=None):
    _box(slide, 0, 0, W, Inches(1.05), fill=NAVY)
    _box(slide, 0, Inches(1.05), W, Pt(4), fill=AMBER)
    _text(slide, Inches(0.5), Inches(0.16), Inches(11.5), Inches(0.8),
          [(f"{idx:02d}  ", 26, True, AMBER), (title, 26, True, WHITE)],
          anchor=MSO_ANCHOR.MIDDLE)
    if pct:
        _text(slide, Inches(11.6), Inches(0.16), Inches(1.5), Inches(0.8),
              [(pct, 16, True, AMBER)], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def _bar_chart(slide, x, y, w, h, cats, series, chart_type=XL_CHART_TYPE.BAR_CLUSTERED,
               colors=None, title=None, legend=False):
    cd = CategoryChartData()
    cd.categories = cats
    for name, vals in series:
        cd.add_series(name, vals)
    gf = slide.shapes.add_chart(chart_type, x, y, w, h, cd)
    ch = gf.chart
    ch.has_title = bool(title)
    if title:
        ch.chart_title.text_frame.text = title
        ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
        ch.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = INK
    ch.has_legend = legend
    if legend:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
    for plot in ch.plots:
        plot.gap_width = 60
        for i, s in enumerate(plot.series):
            if colors:
                if isinstance(colors[i], list):
                    for j, pt in enumerate(s.points):
                        pt.format.fill.solid()
                        pt.format.fill.fore_color.rgb = colors[i][j]
                else:
                    s.format.fill.solid()
                    s.format.fill.fore_color.rgb = colors[i]
    try:
        ch.font.size = Pt(10)
    except Exception:
        pass
    return ch


# =============================================================================
def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    lang = D.language_signal()
    db = D.database_signal()
    sal = D.salary_by_language()
    role = D.role_priority()
    ide = D.ide_overall()
    rec = D.strategic_recommendations()
    ov = D.dataset_overview().set_index("metric")["value"]

    # ---- SLIDE 1: Title ----
    s = prs.slides.add_slide(blank)
    _bg(s, NAVY)
    _box(s, 0, Inches(3.05), W, Pt(5), fill=AMBER)
    _text(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.4),
          [("A Day in the Life of a Data Analyst", 44, True, WHITE)])
    _text(s, Inches(0.8), Inches(3.25), Inches(11.7), Inches(1.0),
          [("Phân tích Kỹ năng Công nghệ Thông tin đang được thị trường yêu cầu nhất", 24, False, RGBColor(0xCB, 0xD5, 0xE1))])
    _text(s, Inches(0.8), Inches(4.6), Inches(11.7), Inches(1.6),
          [[("VTI Academy", 18, True, AMBER)],
           [("Chuyên viên Phân tích Dữ liệu  ·  Công cụ: Python · SQL · Power BI", 16, False, WHITE)],
           [("Nguồn dữ liệu: Stack Overflow Survey (11.552 dev / 135 quốc gia) + ItViec", 14, False, RGBColor(0x94, 0xA3, 0xB8))]])
    _notes(s, "Xin chào hội đồng. Em trình bày dự án 'A Day in the Life of a Data Analyst'. "
              "Vào vai Chuyên viên Phân tích Dữ liệu tại VTI Academy, em phân tích kỹ năng IT thị trường "
              "đang cần nhất để giúp trung tâm cập nhật chương trình đào tạo. Câu hỏi cốt lõi: "
              "Đào tạo gì để học viên ra trường có việc ngay?")

    # ---- SLIDE 2: Method ----
    s = prs.slides.add_slide(blank)
    _bg(s, LIGHT)
    _header(s, 2, "Phương pháp: Thu thập & Làm sạch dữ liệu", "20%")
    # 3 nguồn
    src = [("Stack Overflow Survey", "CSV · 11.552 dev · 90 cột · 135 quốc gia", BLUE),
           ("ItViec Job Postings", "Web scraping (Playwright) · tin tuyển dụng IT VN", GREEN),
           ("Training Portals", "API khóa học (Coursera/Udemy) — đối chiếu cung đào tạo", AMBER)]
    for i, (t, d, c) in enumerate(src):
        x = Inches(0.5 + i * 4.25)
        _box(s, x, Inches(1.5), Inches(4.0), Inches(1.7), fill=WHITE, line=c)
        _box(s, x, Inches(1.5), Inches(0.12), Inches(1.7), fill=c)
        _text(s, x + Inches(0.25), Inches(1.65), Inches(3.6), Inches(1.4),
              [[(t, 16, True, INK)], [("", 6, False, INK)], [(d, 12, False, GREY)]])
    # pipeline
    _text(s, Inches(0.5), Inches(3.5), Inches(12), Inches(0.5),
          [("Pipeline kỹ thuật", 18, True, NAVY)])
    steps = ["Scrape\n(Playwright)", "Wrangle\n(Pandas)", "Store\n(SQLite + SQL)", "Visualize\n(Power BI)"]
    for i, st in enumerate(steps):
        x = Inches(0.5 + i * 3.15)
        _box(s, x, Inches(4.1), Inches(2.7), Inches(1.0), fill=NAVY)
        _text(s, x, Inches(4.2), Inches(2.7), Inches(0.85),
              [[(st.split(chr(10))[0], 15, True, WHITE)], [(st.split(chr(10))[1], 11, False, RGBColor(0xCB,0xD5,0xE1))]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < 3:
            _text(s, x + Inches(2.7), Inches(4.2), Inches(0.45), Inches(0.85),
                  [("→", 26, True, AMBER)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _box(s, Inches(0.5), Inches(5.45), Inches(12.3), Inches(1.1), fill=WHITE, line=GREEN)
    _text(s, Inches(0.7), Inches(5.55), Inches(12), Inches(0.95),
          [[("Tự động hóa: ", 14, True, GREEN), ("GitHub Actions chạy cron 07:00 mỗi ngày — tự cào ItViec, "
            "kiểm tra 5 quality-check, commit dữ liệu ngược repo. ", 13, False, INK)],
           [("Làm sạch: ", 14, True, GREEN), ("đổi tên cột → snake_case, khử trùng lặp, điền thiếu (categorical→Unknown, "
            "numeric→median), chuẩn hóa tên skill qua từ điển alias, winsorize lương để cắt ngoại lai.", 13, False, INK)]])
    _notes(s, "Em thu thập từ 3 nguồn để đảm bảo đa dạng theo đề bài. Điểm khác biệt: em không chỉ làm 1 lần "
              "mà xây hệ thống tự động trên GitHub Actions — mỗi sáng 7h tự cào ItViec, làm sạch, kiểm tra chất lượng "
              "và commit kết quả. Dữ liệu luôn cập nhật. Về làm sạch: chuẩn hóa tên cột snake_case, xử lý thiếu, "
              "chuẩn hóa tên skill, và winsorize lương để trung vị không bị méo bởi giá trị bất thường.")

    # ---- SLIDE 3: Insight 1 - Languages ----
    s = prs.slides.add_slide(blank)
    _bg(s, LIGHT)
    _header(s, 3, "Insight 1 — Ngôn ngữ lập trình được yêu cầu nhất")
    top = lang.head(8).iloc[::-1]
    _bar_chart(s, Inches(0.5), Inches(1.4), Inches(7.6), Inches(5.6),
               list(top["language"]), [("Muốn dùng năm tới", list(top["desired_next_year"]))],
               colors=[BLUE], title="Top ngôn ngữ theo nhu cầu năm tới (số dev)")
    _box(s, Inches(8.4), Inches(1.5), Inches(4.5), Inches(5.4), fill=WHITE, line=BLUE)
    _text(s, Inches(8.65), Inches(1.7), Inches(4.1), Inches(5.1),
          [[("Phát hiện chính", 17, True, NAVY)], [("", 6, False, INK)],
           [("JavaScript #1", 15, True, INK), (" — vẫn thống lĩnh web (58% muốn dùng).", 13, False, GREY)],
           [("Python #2", 15, True, GREEN), (" — ngôn ngữ DUY nhất trong top 4 vẫn tăng (+15%).", 13, False, GREY)],
           [("SQL #3", 15, True, INK), (" — kỹ năng nền tảng cho mọi vai trò dữ liệu.", 13, False, GREY)],
           [("", 8, False, INK)],
           [("Hành động HR", 15, True, AMBER)],
           [("Tách rõ kỹ năng nền tảng (JS/SQL/Python) với kỹ năng tăng trưởng khi thiết kế curriculum & JD.", 13, False, INK)]])
    _notes(s, "JavaScript dẫn đầu — không bất ngờ với web. Nhưng điểm quan trọng: trong top 4, chỉ Python còn TĂNG. "
              "JavaScript, SQL, HTML/CSS đều giảm net vì đã bão hòa. SQL #3 khẳng định kỹ năng database là nền tảng. "
              "Bài học HR: phân biệt rõ kỹ năng nền với kỹ năng tăng trưởng.")

    # ---- SLIDE 4: Insight 2 - Database & IDE ----
    s = prs.slides.add_slide(blank)
    _bg(s, LIGHT)
    _header(s, 4, "Insight 2 — Cơ sở dữ liệu & IDE")
    d = db.head(6).iloc[::-1]
    _bar_chart(s, Inches(0.4), Inches(1.4), Inches(6.3), Inches(5.3),
               list(d["database"]),
               [("Đang dùng", list(d["worked"])), ("Muốn dùng", list(d["desired_next_year"]))],
               colors=[GREY, BLUE], title="Database: đang dùng vs muốn dùng", legend=True)
    di = ide.head(6).iloc[::-1]
    _bar_chart(s, Inches(6.9), Inches(1.4), Inches(6.0), Inches(5.3),
               list(di["ide"]), [("Số dev", list(di["developer_count"]))],
               colors=[AMBER], title="IDE phổ biến nhất")
    _text(s, Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.6),
          [[("PostgreSQL", 13, True, GREEN), (" vừa phổ biến vừa còn tăng (+5.6%) → chọn làm DB trục chính.  ", 12, False, INK),
            ("VS Code", 13, True, AMBER), (" chiếm ~55% mọi vai trò → IDE chuẩn onboarding.", 12, False, INK)]])
    _notes(s, "PostgreSQL vượt MySQL ở khảo sát toàn cầu và vẫn tăng trưởng dương, trong khi MySQL giảm 40%. "
              "MongoDB và Redis là lớp NoSQL bổ trợ đang lên. Về IDE, VS Code chiếm ~55% — vượt xa phần còn lại. "
              "Kết luận: VTI nên lấy PostgreSQL và VS Code làm chuẩn chung khi onboarding.")

    # ---- SLIDE 5: Statistical - emerging vs declining ----
    s = prs.slides.add_slide(blank)
    _bg(s, LIGHT)
    _header(s, 5, "Phân tích Thống kê — Emerging vs Declining", "50%")
    order = lang.reindex(lang["net_change"].abs().sort_values().index).tail(10)
    colors = [GREEN if v >= 0 else RED for v in order["net_change"]]
    _bar_chart(s, Inches(0.5), Inches(1.4), Inches(7.6), Inches(5.6),
               list(order["language"]), [("Net change", list(order["net_change"]))],
               colors=[colors], title="Thay đổi ròng nhu cầu (muốn − đang dùng)")
    _box(s, Inches(8.4), Inches(1.5), Inches(4.5), Inches(5.4), fill=WHITE, line=GREEN)
    _text(s, Inches(8.65), Inches(1.65), Inches(4.1), Inches(5.2),
          [[("2 kỹ thuật thống kê", 16, True, NAVY)],
           [("1. Tần suất & tỷ lệ: ", 13, True, INK), ("đếm lượt skill / tổng dev.", 12, False, GREY)],
           [("2. Net-change & growth %: ", 13, True, INK), ("(muốn − đang)/đang.", 12, False, GREY)],
           [("", 6, False, INK)],
           [("Tăng mạnh ", 14, True, GREEN), ("(mở khóa mới):", 13, True, INK)],
           [("Rust +370% · Kotlin +153% · Go +148% · TypeScript +27%", 12, True, GREEN)],
           [("", 4, False, INK)],
           [("Đi xuống ", 14, True, RED), ("(không mở rộng):", 13, True, INK)],
           [("PHP −50% · Java −34% · Bash −33% · HTML/CSS −32%", 12, True, RED)],
           [("", 6, False, INK)],
           [("→ VTI không nên mở thêm khóa PHP; nên mở Rust/Go/TypeScript.", 12, False, INK)]])
    _notes(s, "Đây là phân tích thống kê quan trọng nhất. Em so sánh ngôn ngữ đang dùng với ngôn ngữ muốn dùng năm sau, "
              "tính net-change và growth %. Rust, Kotlin, Go, TypeScript tăng rất mạnh — đây là tương lai. "
              "PHP, Java, Bash đang co lại. Insight chiến lược: không mở thêm PHP, nên đầu tư Rust/Go/TypeScript.")

    # ---- SLIDE 6: Salary x Growth ----
    s = prs.slides.add_slide(blank)
    _bg(s, LIGHT)
    _header(s, 6, "Insight 3 — Lương vs Tăng trưởng → chọn skill 'vàng'")
    sd = sal.head(10).iloc[::-1]
    _bar_chart(s, Inches(0.5), Inches(1.4), Inches(7.6), Inches(5.6),
               list(sd["language"]), [("Median salary USD", list(sd["median_salary"]))],
               colors=[GREEN], title="Lương trung vị theo ngôn ngữ (USD)")
    _box(s, Inches(8.4), Inches(1.5), Inches(4.5), Inches(5.4), fill=WHITE, line=AMBER)
    _text(s, Inches(8.65), Inches(1.7), Inches(4.1), Inches(5.1),
          [[("Giao 3 tín hiệu", 16, True, NAVY), (" = lương + tăng trưởng + nhu cầu", 12, False, GREY)],
           [("", 6, False, INK)],
           [("Go", 14, True, GREEN), (": $80k + tăng 148% → ", 12, False, INK), ("skill vàng", 12, True, GREEN)],
           [("Rust", 14, True, GREEN), (": $70k + tăng 370% → đón đầu", 12, False, INK)],
           [("TypeScript", 14, True, GREEN), (": $60k + tăng 27% → mass-market", 12, False, INK)],
           [("Python", 14, True, BLUE), (": $64k + tăng 15% → nền data/AI", 12, False, INK)],
           [("", 6, False, INK)],
           [("Lưu ý: ", 13, True, RED), ("lương tự khai, có ngoại lai → dùng TRUNG VỊ kèm cỡ mẫu, "
            "winsorize ở mức $300k.", 12, False, GREY)]])
    _notes(s, "Lựa chọn tốt nhất nằm ở giao của 3 tín hiệu: lương tốt, tăng trưởng cao, đủ nhu cầu. "
              "Go nổi bật: lương cao nhất nhóm phổ biến ($80k) và tăng 148%. Rust lương tốt và tăng bùng nổ. "
              "TypeScript là lựa chọn mass-market. Về phương pháp: lương tự khai nên em dùng trung vị kèm cỡ mẫu và winsorize.")

    # ---- SLIDE 7: Dashboard overview ----
    s = prs.slides.add_slide(blank)
    _bg(s, LIGHT)
    _header(s, 7, "Power BI Dashboard — Trang 1: Tổng quan")
    kpis = [("11.552", "Developers"), ("135", "Quốc gia"), ("$57.8k", "Median salary"),
            ("4", "Skill Emerging"), ("PostgreSQL", "DB trục chính")]
    for i, (v, k) in enumerate(kpis):
        x = Inches(0.5 + i * 2.5)
        _box(s, x, Inches(1.35), Inches(2.3), Inches(1.1), fill=WHITE, line=BLUE)
        _text(s, x, Inches(1.45), Inches(2.3), Inches(0.95),
              [[(v, 22, True, NAVY)], [(k, 12, False, GREY)]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    lt = lang.head(7).iloc[::-1]
    _bar_chart(s, Inches(0.5), Inches(2.7), Inches(6.0), Inches(4.2),
               list(lt["language"]), [("Muốn dùng", list(lt["desired_next_year"]))],
               colors=[BLUE], title="Top ngôn ngữ")
    dt = db.head(6).iloc[::-1]
    _bar_chart(s, Inches(6.8), Inches(2.7), Inches(6.0), Inches(4.2),
               list(dt["database"]), [("Muốn dùng", list(dt["desired_next_year"]))],
               colors=[GREEN], title="Top database")
    _text(s, Inches(0.5), Inches(6.95), Inches(12), Inches(0.4),
          [[("Slicers tương tác: ", 12, True, AMBER), ("Country · Signal (Emerging/Stable/Declining) — "
            "5 KPI cards + 2 chart + 2 slicer (đạt yêu cầu tối thiểu).", 11, False, GREY)]])
    _notes(s, "Toàn bộ insight được trực quan hóa trong Power BI. Trang 1 — Overview: 5 thẻ KPI, 2 biểu đồ chính, "
              "2 bộ lọc tương tác (Country, Signal). Đáp ứng đủ yêu cầu tối thiểu của đề bài. "
              "Bản HTML tương tác kèm theo trong powerbi/dashboard.html cho phép demo nhanh không cần cài Power BI.")

    # ---- SLIDE 8: Dashboard deep dive ----
    s = prs.slides.add_slide(blank)
    _bg(s, LIGHT)
    _header(s, 8, "Power BI Dashboard — Trang 2: Deep Dive theo vai trò")
    rp = role.iloc[::-1]
    bands = {"Priority 1": RED, "Priority 2": AMBER, "Priority 3": GREY}
    pcolors = [bands[b] for b in rp["priority_band"]]
    _bar_chart(s, Inches(0.5), Inches(1.4), Inches(7.4), Inches(5.5),
               list(rp["role"]), [("Điểm ưu tiên tuyển dụng", list(rp["hiring_priority_score"]))],
               colors=[pcolors], title="Ưu tiên tuyển dụng theo vai trò (điểm)")
    _box(s, Inches(8.2), Inches(1.5), Inches(4.7), Inches(5.4), fill=WHITE, line=NAVY)
    _text(s, Inches(8.45), Inches(1.7), Inches(4.3), Inches(5.1),
          [[("Đọc dashboard", 16, True, NAVY)],
           [("", 4, False, INK)],
           [("Priority 1 ", 13, True, RED), ("Full-stack — pipeline lõi.", 12, False, INK)],
           [("Priority 2 ", 13, True, AMBER), ("Back-end, DevOps, Front-end.", 12, False, INK)],
           [("Priority 3 ", 13, True, GREY), ("Data, Mobile — tuyển theo nhu cầu.", 12, False, INK)],
           [("", 6, False, INK)],
           [("Tính năng tương tác:", 14, True, AMBER)],
           [("• Slicer Country & Signal", 12, False, INK)],
           [("• Heatmap IDE × vai trò", 12, False, INK)],
           [("• Drill-through skill → công ty ItViec đang tuyển", 12, False, INK)]])
    _notes(s, "Trang 2 — Deep Dive: điểm ưu tiên tuyển dụng quy từ nhu cầu + lương + tăng trưởng. "
              "Full-stack là Priority 1. Heatmap IDE theo vai trò và drill-through từ skill sang danh sách công ty "
              "ItViec đang tuyển biến dashboard thành công cụ định hướng cho học viên.")

    # ---- SLIDE 9: Recommendations ----
    s = prs.slides.add_slide(blank)
    _bg(s, LIGHT)
    _header(s, 9, "Kiến nghị Chiến lược cho VTI Academy", "10%")
    colors = [GREEN, BLUE, NAVY]
    for i, (_, r) in enumerate(rec.iterrows()):
        y = Inches(1.4 + i * 1.85)
        _box(s, Inches(0.5), y, Inches(12.3), Inches(1.65), fill=WHITE, line=colors[i])
        _box(s, Inches(0.5), y, Inches(0.15), Inches(1.65), fill=colors[i])
        _text(s, Inches(0.85), y + Inches(0.12), Inches(11.7), Inches(1.5),
              [[(f"{i+1}. {r['horizon']}: ", 16, True, colors[i]), (r["recommendation"], 15, True, INK)],
               [("Lý do: ", 12, True, GREY), (r["evidence"], 12, False, GREY)]])
    _notes(s, "Ba kiến nghị theo 3 tầm thời gian. Ngắn hạn: mở Cloud/DevOps + Python vì lương cao nhất và cung còn ít. "
              "Trung hạn: đưa TypeScript/Go/Kotlin/Rust vào nâng cao để đón đầu. "
              "Dài hạn: chuẩn hóa nền PostgreSQL/VS Code/Git và duy trì pipeline scraping để tái phân tích mỗi quý — "
              "biến phân tích 1 lần thành năng lực phân tích liên tục.")

    # ---- SLIDE 10: Thanks / Q&A ----
    s = prs.slides.add_slide(blank)
    _bg(s, NAVY)
    _box(s, 0, Inches(3.5), W, Pt(5), fill=AMBER)
    _text(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.2),
          [("Cảm ơn & Q&A", 48, True, WHITE)], align=PP_ALIGN.CENTER)
    _text(s, Inches(0.8), Inches(3.8), Inches(11.7), Inches(1.5),
          [[("Sản phẩm: Clean dataset · Python (scraping + wrangling) · SQL · Power BI · Slides", 16, False, RGBColor(0xCB,0xD5,0xE1))],
           [("GitHub: github.com/quangnm0208/vti-course", 15, True, AMBER)],
           [("Email: quangnm0208@gmail.com", 14, False, WHITE)]], align=PP_ALIGN.CENTER)
    _notes(s, "Em xin kết thúc phần trình bày. Cảm ơn hội đồng đã lắng nghe. Em sẵn sàng trả lời câu hỏi và "
              "demo code Python, file SQL hoặc dashboard trực tiếp trên repo.")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"[ok] {len(prs.slides.__iter__.__self__._sldIdLst)} slide -> {OUT.relative_to(ROOT)}")


if __name__ == "__main__":
    build()
